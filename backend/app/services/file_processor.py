import os
from typing import List
from pathlib import Path
import PyPDF2
from pptx import Presentation
from openpyxl import load_workbook
from docx import Document


class FileProcessor:
    """파일에서 텍스트를 추출하는 클래스"""
    
    @staticmethod
    def extract_text(file_path: str) -> str:
        """파일 확장자에 따라 텍스트 추출"""
        extension = Path(file_path).suffix.lower()
        
        try:
            if extension == '.txt':
                return FileProcessor._extract_from_txt(file_path)
            elif extension == '.pdf':
                return FileProcessor._extract_from_pdf(file_path)
            elif extension in ['.pptx', '.ppt']:
                return FileProcessor._extract_from_pptx(file_path)
            elif extension in ['.xlsx', '.xls']:
                return FileProcessor._extract_from_xlsx(file_path)
            elif extension in ['.docx', '.doc']:
                return FileProcessor._extract_from_docx(file_path)
            else:
                raise ValueError(f"지원하지 않는 파일 형식: {extension}")
        except Exception as e:
            raise Exception(f"파일 처리 중 오류 발생: {str(e)}")
    
    @staticmethod
    def _extract_from_txt(file_path: str) -> str:
        """TXT 파일에서 텍스트 추출"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        """PDF 파일에서 텍스트 추출"""
        text = []
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
        return '\n'.join(text)
    
    @staticmethod
    def _extract_from_pptx(file_path: str) -> str:
        """PPTX 파일에서 텍스트 추출"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    
    @staticmethod
    def _extract_from_xlsx(file_path: str) -> str:
        """
        XLSX 파일에서 구조화된 텍스트 추출
        - 다중 행 헤더 지원
        - 병합된 셀 처리
        - 마크다운 테이블 형식 변환
        """
        wb = load_workbook(file_path, data_only=True)
        all_sheets_text = []
        
        for sheet in wb:
            sheet_name = sheet.title
            all_sheets_text.append(f"\n{'='*80}")
            all_sheets_text.append(f"시트명: {sheet_name}")
            all_sheets_text.append(f"{'='*80}\n")
            
            # 모든 행 데이터 수집
            rows_data = []
            for row in sheet.iter_rows(values_only=True):
                # None이 아닌 셀이 하나라도 있으면 유효한 행
                if any(cell is not None for cell in row):
                    rows_data.append(row)
            
            if not rows_data:
                all_sheets_text.append("(빈 시트)\n")
                continue
            
            # 다중 행 헤더 감지 및 병합
            headers = FileProcessor._extract_multirow_headers(rows_data, sheet)
            
            # 헤더가 유효한지 확인
            has_header = headers and any(h.strip() for h in headers)
            
            if has_header:
                # 헤더 행 개수 결정 (보통 1~2행)
                header_rows_count = FileProcessor._detect_header_rows_count(rows_data)
                data_start_row = header_rows_count
                
                all_sheets_text.append("【테이블 데이터】\n")
                
                # 마크다운 테이블 형식으로 변환
                # 헤더 출력
                all_sheets_text.append("| " + " | ".join(headers) + " |")
                all_sheets_text.append("|" + "|".join(["---" for _ in headers]) + "|")
                
                # 데이터 행 출력 (헤더 행 이후부터)
                for row in rows_data[data_start_row:]:
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    # 헤더 개수에 맞춰 조정
                    while len(row_values) < len(headers):
                        row_values.append("")
                    row_values = row_values[:len(headers)]
                    all_sheets_text.append("| " + " | ".join(row_values) + " |")
                
                all_sheets_text.append("")
                
                # 추가로 검색이 잘 되도록 행별 텍스트도 추가
                all_sheets_text.append("\n【행별 상세 정보】\n")
                for idx, row in enumerate(rows_data[data_start_row:], start=data_start_row+1):
                    row_info = []
                    for header, value in zip(headers, row):
                        if value is not None and str(value).strip():
                            row_info.append(f"{header}: {value}")
                    if row_info:
                        all_sheets_text.append(f"행 {idx}번: " + ", ".join(row_info))
            else:
                # 헤더가 없는 경우 일반 텍스트로
                all_sheets_text.append("【데이터】\n")
                for idx, row in enumerate(rows_data, start=1):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        all_sheets_text.append(f"행 {idx}: " + " | ".join(row_text))
            
            all_sheets_text.append("")
        
        return '\n'.join(all_sheets_text)
    
    @staticmethod
    def _detect_header_rows_count(rows_data: List[tuple]) -> int:
        """
        헤더 행 개수 자동 감지
        - 첫 번째 행에 빈 셀이 많고, 두 번째 행에도 헤더 같은 값이 있으면 2행
        - 그 외에는 1행
        """
        if len(rows_data) < 2:
            return 1
        
        first_row = rows_data[0]
        second_row = rows_data[1]
        
        # 첫 번째 행의 빈 셀 비율
        empty_count_first = sum(1 for cell in first_row if cell is None or str(cell).strip() == "")
        empty_ratio = empty_count_first / len(first_row) if len(first_row) > 0 else 0
        
        # 두 번째 행에 "마일스톤", "목표일" 같은 헤더 패턴이 있는지 확인
        second_row_has_headers = any(
            cell and isinstance(cell, str) and len(str(cell).strip()) > 0 and str(cell).strip() not in ['NaN', 'nan']
            for cell in second_row
        )
        
        # 두 번째 행의 값들이 대부분 비어있지만 일부만 값이 있으면 헤더일 가능성
        non_empty_second = [cell for cell in second_row if cell is not None and str(cell).strip()]
        
        # 2행 헤더 조건: 첫 번째 행에 빈 셀이 20% 이상이고, 두 번째 행에 1~3개 정도의 헤더값
        if empty_ratio > 0.2 and 1 <= len(non_empty_second) <= 3 and second_row_has_headers:
            return 2
        
        return 1
    
    @staticmethod
    def _extract_multirow_headers(rows_data: List[tuple], sheet) -> List[str]:
        """
        다중 행 헤더 추출 및 병합
        - 첫 번째 행과 두 번째 행을 결합하여 완전한 헤더 생성
        - 빈 셀은 두 번째 행 값으로 채움
        """
        if not rows_data:
            return []
        
        first_row = rows_data[0]
        headers = []
        
        # 두 번째 행 존재 여부 확인
        has_second_row = len(rows_data) > 1
        second_row = rows_data[1] if has_second_row else None
        
        for col_idx, cell in enumerate(first_row):
            header = str(cell) if cell is not None and str(cell).strip() else ""
            
            # 첫 번째 행이 비어있거나 의미없는 값이면 두 번째 행 확인
            if not header or header in ['None', 'nan', 'NaN']:
                if second_row and col_idx < len(second_row):
                    second_val = second_row[col_idx]
                    if second_val is not None and str(second_val).strip():
                        header = str(second_val).strip()
                    else:
                        header = f"열{col_idx+1}"
                else:
                    header = f"열{col_idx+1}"
            else:
                # 첫 번째 행에 값이 있어도, 두 번째 행에 추가 정보가 있으면 결합
                if second_row and col_idx < len(second_row):
                    second_val = second_row[col_idx]
                    # 두 번째 행 값이 있고, 첫 번째 행과 다르면 결합
                    if (second_val is not None and 
                        str(second_val).strip() and 
                        str(second_val).strip() != header and
                        str(second_val).strip() not in ['None', 'nan', 'NaN']):
                        # "마일스톤" / "목표일" 같은 경우 처리
                        if header == str(second_val).strip():
                            # 같은 값이면 중복 방지
                            pass
                        else:
                            header = f"{header}_{second_val}".strip('_')
            
            headers.append(header)
        
        return headers
    
    @staticmethod
    def _extract_from_docx(file_path: str) -> str:
        """DOCX 파일에서 텍스트 추출"""
        doc = Document(file_path)
        text = [paragraph.text for paragraph in doc.paragraphs]
        return '\n'.join(text)

